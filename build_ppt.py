"""Generates the FormForge management + IT presentation (.pptx).

Run:   python3 /app/build_ppt.py
Out:   /app/FormForge_Presentation.pptx

The deck is 16:9, uses a modern slate/blue palette, and mixes narrative
sections for management with technical detail sections for IT.
"""
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Palette + helpers
# ---------------------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x1F, 0x3D)   # slate-900
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # blue-600
LIGHT  = RGBColor(0xEF, 0xF6, 0xFF)   # blue-50
GREY   = RGBColor(0x64, 0x74, 0x8B)   # slate-500
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
GREEN  = RGBColor(0x10, 0xB9, 0x81)   # emerald-500
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # amber-500

TITLE_FONT = "Calibri"
BODY_FONT  = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font=BODY_FONT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=16, color=DARK):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # bullet dot
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = BODY_FONT
        dot.font.size = Pt(size)
        dot.font.color.rgb = BLUE
        dot.font.bold = True
        # body
        run = p.add_run()
        run.text = item
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_footer(slide, page):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(6), Inches(0.3),
             "FormForge — Enterprise Form & Document Platform",
             size=9, color=GREY)
    add_text(slide, Inches(11.4), Inches(7.1), Inches(1.5), Inches(0.3),
             f"Page {page}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


def add_header_bar(slide):
    """Thin coloured strip at the top."""
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.15), BLUE)


def title_slide_h(slide, title, subtitle=None):
    add_header_bar(slide)
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.6),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GREY)
    # divider
    add_rect(slide, Inches(0.6), Inches(1.55), Inches(1.2), Emu(30_000), BLUE)


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
# Full navy background
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
# Accent stripe
add_rect(s, 0, Inches(3.6), prs.slide_width, Inches(0.05), BLUE)
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.2),
         "FormForge", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(2.6), Inches(12), Inches(0.6),
         "Enterprise Form & Document Management Platform",
         size=22, color=LIGHT)
add_text(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.4),
         "For Top Management & IT Leadership",
         size=14, bold=True, color=BLUE)
add_text(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.4),
         "Standard Forms   •   PDF Forms   •   Plant Documents   •   Backup & Restore",
         size=13, color=LIGHT)
add_text(s, Inches(0.7), Inches(6.6), Inches(6), Inches(0.4),
         "Prepared: " + datetime.now().strftime("%B %Y"),
         size=11, color=GREY)
add_text(s, Inches(9), Inches(6.6), Inches(4), Inches(0.4),
         "Version 1.0", size=11, color=GREY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Slide 2 — Executive Summary
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Executive Summary",
              "The problem, the platform, and the payoff")
add_text(s, Inches(0.6), Inches(1.9), Inches(12.2), Inches(0.5),
         "Problem", size=16, bold=True, color=BLUE)
add_text(s, Inches(0.6), Inches(2.35), Inches(12.2), Inches(0.7),
         "Field data collection, plant document control, and multi-tier approvals "
         "were fragmented across Excel, email, and shared drives — creating "
         "compliance risk, poor visibility, and slow decision cycles.",
         size=13, color=DARK)

add_text(s, Inches(0.6), Inches(3.35), Inches(12.2), Inches(0.5),
         "Solution", size=16, bold=True, color=BLUE)
add_bullets(s, Inches(0.9), Inches(3.8), Inches(12), Inches(1.6), [
    "Single web platform for Standard & PDF forms, plant document vaults, and approvals.",
    "Strict 4-tier RBAC (Super Admin › Region › Cluster › Vendor) with row-level scoping.",
    "One-click disaster-recovery snapshots + migration bundles + upload-restore.",
], size=14)

add_text(s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(0.5),
         "Impact", size=16, bold=True, color=BLUE)
add_bullets(s, Inches(0.9), Inches(5.85), Inches(12), Inches(1.2), [
    "Real-time dashboards across every submission, region, and plant.",
    "Zero-loss disaster recovery — full restore in under 5 minutes.",
    "Ready for on-prem or cloud; single-command Docker deploy.",
], size=14)
add_footer(s, 2)


# ---------------------------------------------------------------------------
# Slide 3 — What FormForge Does
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "What FormForge Does",
              "Four modules under one roof")

# 4-column layout
COLS = [
    ("Form Builder", "Drag-and-drop editor for standard forms with 20+ field types, conditional logic, plant lookups, and public URLs.", BLUE),
    ("PDF Forms", "Overlay text, checkboxes, radios and signatures onto existing PDF templates. Downloads keep native PDF layout.", GREEN),
    ("Plant Documents", "Per-plant document vaults with folder templates, drag-and-drop upload, inline preview, and zip-download.", AMBER),
    ("Backup & Migration", "One-click MongoDB + file snapshots. Restore from server-side or an uploaded .tar.gz. Move stacks in minutes.", NAVY),
]
col_w = Inches(2.85)
gap = Inches(0.2)
start_x = Inches(0.55)
for i, (title, body, colour) in enumerate(COLS):
    x = start_x + i * (col_w + gap)
    y = Inches(2.1)
    h = Inches(4.6)
    add_rect(s, x, y, col_w, h, LIGHT)
    add_rect(s, x, y, col_w, Inches(0.35), colour)
    add_text(s, x + Inches(0.2), y + Inches(0.55), col_w - Inches(0.4), Inches(0.5),
             title, size=17, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(1.15), col_w - Inches(0.4), h - Inches(1.3),
             body, size=12, color=DARK)
add_footer(s, 3)


# ---------------------------------------------------------------------------
# Slide 4 — Business Value  (management-focused)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Business Value",
              "Why this matters to leadership")

BENEFITS = [
    ("Compliance-ready", "Every submission is timestamped, audit-logged, and immutable. Full trail for auditors."),
    ("Real-time visibility", "Dashboards roll up submissions across regions, plants, forms, and users."),
    ("Zero-loss operations", "Automatic daily backups + one-click restore = disaster recovery in under 5 min."),
    ("Vendor accountability", "Multi-tier approval workflows enforce region → cluster → vendor sign-offs before actions."),
    ("Faster field cycles", "Public form URLs let vendors submit inspections from mobile without app installs."),
    ("Cost-controlled scaling", "Self-hosted on Docker. Predictable footprint, no per-submission fees."),
]
col_w = Inches(6.1)
gap_y = Inches(0.15)
x0 = Inches(0.55)
y0 = Inches(2.0)
row_h = Inches(1.5)
for i, (t, b) in enumerate(BENEFITS):
    col = i % 2
    row = i // 2
    x = x0 + col * (col_w + Inches(0.15))
    y = y0 + row * (row_h + gap_y)
    add_rect(s, x, y, col_w, row_h, WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(0.1), row_h, BLUE)
    add_text(s, x + Inches(0.3), y + Inches(0.15), col_w - Inches(0.5), Inches(0.45),
             t, size=15, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(0.65), col_w - Inches(0.5), Inches(0.85),
             b, size=11, color=DARK)
add_footer(s, 4)


# ---------------------------------------------------------------------------
# Slide 5 — Feature Snapshot (What Ships Today)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Feature Snapshot",
              "Everything that ships in the current release")

FEATURES = [
    "Standard form builder — 20+ field types, conditional logic, dynamic dropdowns, plant / asset lookups",
    "PDF form builder — per-option checkbox & radio positioning, inline signature, filename templates",
    "Public forms — login-gated with submitter capture, customisable download filenames",
    "Plant documents — auto-provisioned folder templates, rename/delete/zip-download, drag-and-drop upload, inline PDF/image/DOCX viewer",
    "Schedule vs Actual maintenance tracker with mandatory evidence uploads",
    "Master data manager (vendors, plants, custom lookup tables) with CSV import",
    "Multi-tier approvals — deactivation & permission changes route through Region → Cluster → Vendor Admin",
    "Backup & Restore — snapshot download, migration bundle, restore-from-file",
    "Dashboard — trend charts, per-form analytics, storage usage, recent activity (includes PDF submissions)",
    "SMTP email notifications, in-app notifications, exportable audit logs",
]
add_bullets(s, Inches(0.7), Inches(2.0), Inches(12), Inches(5.0), FEATURES, size=13)
add_footer(s, 5)


# ---------------------------------------------------------------------------
# Slide 6 — RBAC & Security  (IT-focused)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Role-Based Access & Security",
              "Four tiers, strict row-level scoping")

# Left: RBAC pyramid table
add_text(s, Inches(0.55), Inches(1.9), Inches(6), Inches(0.4),
         "4-tier hierarchy", size=15, bold=True, color=BLUE)

ROLES = [
    ("Super Admin", "Global. Configure system, users, backups.", NAVY),
    ("Admin (Region)", "See only their region's plants, submissions, users.", BLUE),
    ("Cluster Admin", "Owns cluster of vendors within a region.", GREEN),
    ("Vendor Admin", "See own vendor's users + sites only.", AMBER),
    ("Vendor User", "Only assigned sites; cannot see other vendor plants.", GREY),
]
y = Inches(2.4)
for name, desc, colour in ROLES:
    add_rect(s, Inches(0.55), y, Inches(6.1), Inches(0.75), WHITE,
             line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, Inches(0.55), y, Inches(0.14), Inches(0.75), colour)
    add_text(s, Inches(0.85), y + Inches(0.05), Inches(2.2), Inches(0.35),
             name, size=13, bold=True, color=NAVY)
    add_text(s, Inches(0.85), y + Inches(0.4), Inches(5.6), Inches(0.35),
             desc, size=10, color=DARK)
    y += Inches(0.85)

# Right: Security controls
add_text(s, Inches(7.0), Inches(1.9), Inches(6), Inches(0.4),
         "Security controls", size=15, bold=True, color=BLUE)
CTRL = [
    "bcrypt password hashing (12 rounds)",
    "JWT access tokens (32-byte secret enforced in strict mode)",
    "Row-level filters at the query layer — not just UI",
    "Tiered approval workflow for privileged actions",
    "Every mutation writes an immutable audit-log entry",
    "CORS lockdown + optional HTTPS via nginx TLS block",
    "Backups gated to super_admin only",
    "Public forms require login unless explicitly opened",
]
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.8), Inches(4.5), CTRL, size=12)
add_footer(s, 6)


# ---------------------------------------------------------------------------
# Slide 7 — Architecture
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Architecture",
              "React SPA · FastAPI · MongoDB — all containerised")

# Draw 3 stacked boxes for tiers
def tier_box(x, y, w, h, title, sub, colour):
    add_rect(s, x, y, w, h, colour)
    add_text(s, x, y + Inches(0.15), w, Inches(0.4), title,
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.55), w, Inches(0.35), sub,
             size=10, color=LIGHT, align=PP_ALIGN.CENTER)

x = Inches(0.6); w = Inches(3.6); h = Inches(1.0)
tier_box(x,           Inches(2.0), w, h, "React SPA",           "Vite/Craco · Tailwind · shadcn/ui", BLUE)
tier_box(x + Inches(4.4), Inches(2.0), w, h, "FastAPI Backend", "Python 3.11 · async · JWT + RBAC", NAVY)
tier_box(x + Inches(8.8), Inches(2.0), w, h, "MongoDB",         "20+ collections · mongodump-backed", DARK)

# Connectors (rectangles for arrows)
def arrow(x1, y1, x2, y2):
    ln = s.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = GREY
    ln.line.width = Pt(1.5)
arrow(Inches(4.2), Inches(2.5), Inches(5.0), Inches(2.5))
arrow(Inches(8.6), Inches(2.5), Inches(9.4), Inches(2.5))

# Support systems row
add_text(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.4),
         "Support systems", size=14, bold=True, color=BLUE)

SUPP = [
    ("nginx gateway", "TLS termination · WS · /api → 8001 · /* → 3000"),
    ("Docker Compose", "4 services · named volumes · single-command boot"),
    ("Local file storage", "Configurable roots via .env — no S3 dependency"),
    ("mongodump / mongorestore", "Full DB dump into portable .archive.gz"),
    ("SMTP", "Optional — for approvals & notifications"),
    ("mammoth.js", "Client-side .docx → HTML for inline preview"),
]
col_w = Inches(4.05); gap = Inches(0.1); x0 = Inches(0.55); y0 = Inches(3.9)
for i, (n, d) in enumerate(SUPP):
    row = i // 3; col = i % 3
    x = x0 + col * (col_w + gap)
    y = y0 + row * Inches(1.05)
    add_rect(s, x, y, col_w, Inches(0.95), LIGHT)
    add_text(s, x + Inches(0.2), y + Inches(0.1), col_w - Inches(0.3), Inches(0.35),
             n, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), y + Inches(0.5), col_w - Inches(0.3), Inches(0.4),
             d, size=10, color=DARK)
add_footer(s, 7)


# ---------------------------------------------------------------------------
# Slide 8 — Disaster Recovery
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Disaster Recovery Story",
              "Fully self-serve; no vendor lock-in")

STEPS = [
    ("1. Daily auto-snapshot",  "Automatic tarball at configurable UTC time — MongoDB dump + all upload roots. Retention window is configurable (default 3 days).", BLUE),
    ("2. Manual snapshot",       "Super-admin can trigger anytime from Settings → Backup & Restore. Takes ~30 seconds for a mid-size deployment.", GREEN),
    ("3. Migration bundle",      "Single button generates and downloads a portable .tar.gz that can be restored on any host via ./migrate.sh or the in-app uploader.", AMBER),
    ("4. Restore from file",     "Upload a bundle from disk; the app restores Mongo + uploads. Total time from click to running system: 30 sec – 5 min depending on size.", NAVY),
]
y = Inches(2.0)
for title, body, colour in STEPS:
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.15), LIGHT)
    add_rect(s, Inches(0.6), y, Inches(0.15), Inches(1.15), colour)
    add_text(s, Inches(0.95), y + Inches(0.15), Inches(11), Inches(0.35),
             title, size=15, bold=True, color=NAVY)
    add_text(s, Inches(0.95), y + Inches(0.55), Inches(11), Inches(0.55),
             body, size=11, color=DARK)
    y += Inches(1.25)
add_footer(s, 8)


# ---------------------------------------------------------------------------
# Slide 9 — Deployment Options
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Deployment Options",
              "Same codebase, three delivery models")

OPTIONS = [
    ("On-Premises Docker",
     "One-command boot via docker compose.  4 services (Mongo + backend + frontend + gateway).  Ideal for HSE, plant-ops, restricted-network deployments.",
     ["No egress required after image build", "Full data sovereignty", "Runs on Windows Docker Desktop or Linux"]),
    ("Cloud VM",
     "Same Docker stack on any VPS (AWS EC2, Azure VM, GCP Compute).  nginx TLS block + certbot for HTTPS.",
     ["Elastic scale-up", "Managed backups to S3-compatible storage", "Public URL / custom domain"]),
    ("Local (Dev / Field-Rugged)",
     "Native Windows / Linux run — Python + Node + MongoDB Community.  Perfect for offline field boxes or air-gapped labs.",
     ["No container overhead", "Native MongoDB installer", "MONGODUMP_BIN env-var for portable backups"]),
]
col_w = Inches(4.05); gap = Inches(0.15); x0 = Inches(0.55); y0 = Inches(2.0)
for i, (name, blurb, feats) in enumerate(OPTIONS):
    x = x0 + i * (col_w + gap)
    add_rect(s, x, y0, col_w, Inches(4.9), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y0, col_w, Inches(0.5), BLUE)
    add_text(s, x, y0 + Inches(0.09), col_w, Inches(0.4),
             name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), y0 + Inches(0.65), col_w - Inches(0.4), Inches(1.5),
             blurb, size=11, color=DARK)
    add_bullets(s, x + Inches(0.2), y0 + Inches(2.4), col_w - Inches(0.4), Inches(2.2),
                feats, size=11)
add_footer(s, 9)


# ---------------------------------------------------------------------------
# Slide 10 — Tech Stack (IT)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Technology Stack",
              "Modern, well-supported, no proprietary dependencies")

STACK = [
    ("Frontend",   ["React 18 + Create React App", "Tailwind CSS + shadcn/ui", "React-PDF (pdf.js)", "mammoth.js (.docx render)", "Axios · React-Router · Recharts"]),
    ("Backend",    ["Python 3.11", "FastAPI + Uvicorn", "Motor (async MongoDB)", "PyMuPDF + pypdf (PDF ops)", "PyJWT · bcrypt · Pydantic v2"]),
    ("Database",   ["MongoDB 7", "20+ collections", "mongodump / mongorestore", "Indexes on user_id, form_id, site_id"]),
    ("Infra",      ["Docker · Docker Compose", "nginx (gateway + static)", "Named volumes for uploads", "Optional SMTP relay"]),
    ("Tooling",    ["Yarn (frontend)", "pip freeze pinned reqs", "supervisord in dev", "GitHub Actions ready"]),
    ("Security",   ["JWT auth (256-bit secret)", "bcrypt (12 rounds)", "Query-layer RLS", "Immutable audit log", "CORS lockdown"]),
]
col_w = Inches(4.05); row_h = Inches(2.4); gap = Inches(0.15); x0 = Inches(0.55); y0 = Inches(1.95)
for i, (name, items) in enumerate(STACK):
    col = i % 3; row = i // 3
    x = x0 + col * (col_w + gap); y = y0 + row * (row_h + Inches(0.1))
    add_rect(s, x, y, col_w, row_h, LIGHT)
    add_text(s, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.3), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_bullets(s, x + Inches(0.2), y + Inches(0.6), col_w - Inches(0.3), row_h - Inches(0.7),
                items, size=10)
add_footer(s, 10)


# ---------------------------------------------------------------------------
# Slide 11 — Roadmap
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
title_slide_h(s, "Roadmap",
              "What's next on the platform")

LANES = [
    ("Now — Shipping",
     [
        "Standard + PDF form builders",
        "Plant documents vault (rename / zip / inline viewer)",
        "Full backup / restore / migration bundle",
        "Restore from any local .tar.gz file",
        "Multi-tier approval workflows",
        "Dashboard incl. PDF submissions",
     ], GREEN),
    ("Next — In Progress",
     [
        "Signature field enhancements (pen colour, type-to-sign, save to profile)",
        "Dashboard system-health widget (snapshot age, disk usage, doc counts)",
        "Backup verification endpoint",
        "Bulk file actions (multi-select + zip / delete)",
     ], BLUE),
    ("Later — Backlog",
     [
        "Excel/xlsx inline preview",
        "Push notifications (Web Push / mobile)",
        "SAML / SSO integration",
        "Regional multi-tenant partitioning",
        "S3-compatible remote backup target",
     ], AMBER),
]
col_w = Inches(4.05); x0 = Inches(0.55); y0 = Inches(2.0)
for i, (title, items, colour) in enumerate(LANES):
    x = x0 + i * (col_w + Inches(0.15))
    add_rect(s, x, y0, col_w, Inches(4.9), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y0, col_w, Inches(0.5), colour)
    add_text(s, x, y0 + Inches(0.09), col_w, Inches(0.4),
             title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.2), y0 + Inches(0.7), col_w - Inches(0.4), Inches(4.2),
                items, size=11)
add_footer(s, 11)


# ---------------------------------------------------------------------------
# Slide 12 — Q&A / Contact
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
add_rect(s, 0, Inches(3.4), prs.slide_width, Inches(0.05), BLUE)
add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.2),
         "Questions?", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.7), Inches(3.7), Inches(12), Inches(0.6),
         "Live demo · Deep-dive by module · Roadmap prioritisation",
         size=18, color=LIGHT)
add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
         "Thank you.", size=14, color=BLUE)


# ---------------------------------------------------------------------------
OUT = "/app/FormForge_Presentation.pptx"
prs.save(OUT)
print(f"Wrote {OUT} — {len(prs.slides)} slides")
